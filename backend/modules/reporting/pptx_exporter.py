"""
PowerPoint Rapor Export Modülü
python-pptx kütüphanesi kullanarak PPTX dosyası oluşturur
"""

import logging
import os
from datetime import datetime
from typing import Dict, List, Optional, Tuple

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    logging.info("UYARI: python-pptx kütüphanesi yüklü değil. PowerPoint export çalışmayacak.")
    logging.info("Yüklemek için: pip install python-pptx")

class PowerPointExporter:
    """PowerPoint rapor oluşturucu"""

    def __init__(self) -> None:
        self.prs: Optional[Presentation] = None
        if not PPTX_AVAILABLE:
            raise ImportError("python-pptx kütüphanesi gerekli!")

    def create_presentation(self, title: str = "SUSTAINAGE SDG Raporu") -> Presentation:
        """Yeni sunum oluştur"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(10)
        self.prs.slide_height = Inches(7.5)
        return self.prs

    def add_title_slide(self, title: str, subtitle: str = "") -> None:
        """Başlık slaytı ekle"""
        if self.prs is None:
            self.create_presentation()
        prs = self.prs
        assert prs is not None
        slide_layout = prs.slide_layouts[0]  # Title Slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Başlık
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(44)
        title_shape.text_frame.paragraphs[0].font.bold = True
        title_shape.text_frame.paragraphs[0].font.color.rgb = RGBColor(46, 117, 182)

        # Alt başlık
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            subtitle_shape.text_frame.paragraphs[0].font.size = Pt(28)

        # Tarih ekle
        date_text = f"Oluşturulma Tarihi: {datetime.now().strftime('%d.%m.%Y')}"
        txBox = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = date_text
        p.font.size = Pt(14)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, title: str, content: List[str]) -> None:
        """İçerik slaytı ekle"""
        if self.prs is None:
            self.create_presentation()
        prs = self.prs
        assert prs is not None
        slide_layout = prs.slide_layouts[1]  # Title and Content
        slide = prs.slides.add_slide(slide_layout)

        # Başlık
        title_shape = slide.shapes.title
        title_shape.text = title
        title_shape.text_frame.paragraphs[0].font.size = Pt(32)
        title_shape.text_frame.paragraphs[0].font.bold = True

        # İçerik
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()

        for item in content:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0
            p.font.size = Pt(18)

    def add_table_slide(self, title: str, headers: List[str], data: List[List]) -> None:
        """Tablo slaytı ekle"""
        if self.prs is None:
            self.create_presentation()
        prs = self.prs
        assert prs is not None
        slide_layout = prs.slide_layouts[5]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        # Başlık ekle
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Tablo oluştur
        rows = len(data) + 1  # +1 for headers
        cols = len(headers)

        left = Inches(0.5)
        top = Inches(1.5)
        width = Inches(9)
        height = Inches(5)

        table = slide.shapes.add_table(rows, cols, left, top, width, height).table

        # Başlıkları ekle
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.size = Pt(14)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(46, 117, 182)
            cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

        # Verileri ekle
        for row_idx, row_data in enumerate(data, start=1):
            for col_idx, cell_data in enumerate(row_data):
                cell = table.cell(row_idx, col_idx)
                cell.text = str(cell_data)
                cell.text_frame.paragraphs[0].font.size = Pt(12)

    def add_chart_slide(self, title: str, chart_data: Dict) -> None:
        """Grafik slaytı ekle (basitleştirilmiş)"""
        if self.prs is None:
            self.create_presentation()
        prs = self.prs
        assert prs is not None
        slide_layout = prs.slide_layouts[5]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Başlık
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(32)
        p.font.bold = True
        p.alignment = PP_ALIGN.CENTER

        # Grafik bilgisi (metin olarak - gerçek grafik için daha fazla kod gerekir)
        info_text = "Grafik verileri:\n"
        for key, value in chart_data.items():
            info_text += f"• {key}: {value}\n"

        txBox2 = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf2 = txBox2.text_frame
        tf2.text = info_text
        tf2.paragraphs[0].font.size = Pt(18)

    def save(self, filename: str) -> str:
        """Sunumu kaydet"""
        if self.prs is None:
            raise ValueError("Sunum oluşturulmamış!")

        # Klasör yoksa oluştur
        os.makedirs(os.path.dirname(filename), exist_ok=True)

        self.prs.save(filename)
        return filename

def create_sample_report() -> Tuple[Optional[str], str]:
    try:
        if not PPTX_AVAILABLE:
            return None, "python-pptx kütüphanesi yüklü değil!"

        exporter = PowerPointExporter()
        exporter.create_presentation("SUSTAINAGE SDG Raporu")

        exporter.add_title_slide(
            "SUSTAINAGE SDG",
            "Sürdürülebilir Kalkınma Raporu"
        )

        exporter.add_content_slide(
            "SDG Hedefleri",
            [
                "17 Sürdürülebilir Kalkınma Hedefi",
                "169 Alt Hedef",
                "232 Gösterge",
                "Global standartlara uygun raporlama"
            ]
        )

        exporter.add_table_slide(
            "SDG İlerleme Özeti",
            ["Hedef", "Durum", "İlerleme"],
            [
                ["SDG 1: Yoksulluğa Son", "Devam Ediyor", "%75"],
                ["SDG 7: Temiz Enerji", "Tamamlandı", "%100"],
                ["SDG 13: İklim Eylemi", "Başladı", "%45"],
            ]
        )

        filename = f"reports/SDG_Rapor_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        saved_path = exporter.save(filename)

        return saved_path, "Başarılı!"

    except Exception as e:
        return None, f"Hata: {e}"


def create_kpi_snapshot_presentation(snapshot: Dict, output_dir: str = "reports/kpi_snapshots") -> Tuple[Optional[str], str]:
    try:
        if not PPTX_AVAILABLE:
            return None, "python-pptx kütüphanesi yüklü değil!"

        exporter = PowerPointExporter()
        exporter.create_presentation("SUSTAINAGE KPI Snapshot")

        company = snapshot.get("company") or {}
        period = snapshot.get("period") or {}
        company_name = company.get("name") or "Bilinmeyen Şirket"
        period_label = period.get("label") or period.get("year") or ""
        subtitle = f"{company_name} - {period_label}" if period_label else company_name

        exporter.add_title_slide(
            "KPI Snapshot",
            subtitle
        )

        kpis = snapshot.get("kpis") or []
        if kpis:
            headers = ["Kod", "Ad", "Modül", "Değer", "Birim", "Yıl"]
            rows: List[List] = []
            for kpi in kpis[:100]:
                rows.append(
                    [
                        kpi.get("code") or "",
                        kpi.get("name") or "",
                        kpi.get("module") or "",
                        kpi.get("value"),
                        kpi.get("unit") or "",
                        kpi.get("year") or "",
                    ]
                )
            exporter.add_table_slide("KPI Özeti (İlk 100)", headers, rows)

        os.makedirs(output_dir, exist_ok=True)
        filename = f"kpi_snapshot_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
        full_path = os.path.join(output_dir, filename)
        saved_path = exporter.save(full_path)
        return saved_path, "Başarılı!"

    except Exception as e:
        return None, f"Hata: {e}"


if __name__ == '__main__':
    path, msg = create_sample_report()
    if path:
        logging.info(f"Rapor oluşturuldu: {path}")
    else:
        logging.error(f"Hata: {msg}")


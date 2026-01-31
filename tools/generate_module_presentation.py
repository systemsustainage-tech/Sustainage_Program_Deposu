#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SDG Program Modülleri Sunumu Oluşturucu
- PowerPoint (.pptx) çıktı üretir, resimler ekler
- İçerik: modüller, butonlar/bölümler, raporlama
"""
from __future__ import annotations

import logging
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

OUTPUT_PATH = Path("data/exports/SDG_Modul_Tanitim.pptx")
# Önce resimler/sunum dizinini kullan, yoksa resimler’e düş
IMAGES_DIR = Path("resimler/sunum") if Path("resimler/sunum").exists() else Path("resimler")

# Kurumsal renk paleti
PALETTE = {
    "Kirmizi": "#E53935",
    "Yesil": "#43A047",
    "Mavi": "#1E88E5",
    "Sari": "#FBC02D",
    "Mor": "#8E24AA",
    "Siyah": "#212121",
    "Turkuaz": "#26C6DA",
}

MODULE_COLORS = {
    "SDG": PALETTE["Mavi"],
    "GRI": PALETTE["Mor"],
    "TSRS": PALETTE["Turkuaz"],
    "Karbon": PALETTE["Yesil"],
    "Eslesme": PALETTE["Sari"],
    "Raporlama": PALETTE["Kirmizi"],
    "Yonetim": PALETTE["Siyah"],
}


def hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    return RGBColor(r, g, b)


def is_dark(hex_str: str) -> bool:
    h = hex_str.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    # Perceived luminance
    luminance = 0.299 * r + 0.587 * g + 0.114 * b
    return luminance < 140


def add_title_slide(prs: Presentation, title: str, subtitle: str, image_name: str | None = None) -> None:
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle
    if image_name:
        img_path = IMAGES_DIR / image_name
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(8.0), Inches(1.2), height=Inches(3.5))
    # Üst şerit: palet bloklarını göster
    x = Inches(0.3)
    y = Inches(0.2)
    w = Inches(1.2)
    h = Inches(0.3)
    for i, color in enumerate(PALETTE.values()):
        shape = slide.shapes.add_shape(1, x + Inches(1.25 * i), y, w, h)  # 1: MSO_AUTO_SHAPE_TYPE.RECTANGLE
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(color)
        shape.line.fill.background()


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], image_name: str | None = None, color_hex: str | None = None) -> None:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = b
        p.level = 0
        if color_hex:
            for run in p.runs:
                run.font.color.rgb = hex_to_rgb(color_hex)
    if image_name:
        img_path = IMAGES_DIR / image_name
        if img_path.exists():
            slide.shapes.add_picture(str(img_path), Inches(8.0), Inches(1.5), height=Inches(3.0))
    # Başlık arka planını modül rengine boya
    if color_hex:
        title_shape = slide.shapes.title
        if title_shape:
            title_shape.fill.solid()
            title_shape.fill.fore_color.rgb = hex_to_rgb(color_hex)
            title_shape.line.fill.background()
            # Kontrast için başlık yazı rengi
            title_tf = title_shape.text_frame
            for p in title_tf.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = RGBColor(255, 255, 255) if is_dark(color_hex) else RGBColor(0, 0, 0)


def build_presentation() -> None:
    prs = Presentation()
    add_title_slide(prs, "SDG Program Modülleri", "Modüller, butonlar/bölümler ve raporlama", "main.png")

    # SDG
    sdg_bullets = [
        "Amaç: 17 hedef, 169 alt hedef ve 232 gösterge yönetimi",
        "Bölümler: Mevcut Hedefler, Seçilen Hedefler, İstatistik kartları",
        "Butonlar: İlerleme Takibi, Veri Toplama, Gelişmiş Analiz, Soru Bankası",
        "Butonlar: Veri Doğrulama, Raporlama, Detay Tablosu",
        "Raporlar: SDG ilerleme ve gösterge yanıt raporları",
    ]
    add_bullet_slide(prs, "SDG", sdg_bullets, "SDGs1.jpeg", MODULE_COLORS["SDG"])

    # GRI
    gri_bullets = [
        "Amaç: Global Reporting Initiative standartlarını inceleme ve yanıt",
        "Bölümler: Kategoriler (Universal/Economic/Environmental/Social), İçerik alanı",
        "Butonlar: SDG Seçimine Göre Filtrele, Filtreleri Temizle, CSV Dışa Aktar",
        "Fonksiyonlar: Standart/indikator kartları, yanıt kaydı",
        "Raporlar: CSV dışa aktarma ve GRI uyumlu kayıtlar",
    ]
    add_bullet_slide(prs, "GRI", gri_bullets, None, MODULE_COLORS["GRI"])

    # TSRS
    tsrs_bullets = [
        "Amaç: TSRS standartları ve göstergeleri ile raporlama",
        "Bölümler: Yönetişim, Strateji, Risk Yönetimi, Metrikler",
        "Sekmeler: TSRS Standartları, Raporlama, GRI-SDG Entegrasyon",
        "Butonlar: Kategori butonları (Yönetişim/Strateji/Risk/Metrikler)",
        "Raporlar: TSRSReportingGUI ile TSRS raporlaması",
    ]
    add_bullet_slide(prs, "TSRS", tsrs_bullets, "ESG1.png", MODULE_COLORS["TSRS"])

    # Karbon
    carbon_bullets = [
        "Amaç: Scope 1/2/3 emisyon veri girişi ve hesaplama",
        "Sekmeler: Scope1 (Sabit/Mobil/Kaçak), Scope2 (Elektrik/Isıtma), Scope3",
        "Sekmeler: Hedefler, Azaltma Girişimleri, Raporlar",
        "Butonlar: Kaydet ve Hesapla, CSV dışa aktarma (gerekli yerlerde)",
        "Raporlar: Özet rapor oluşturma ve DB’ye kaydetme",
    ]
    add_bullet_slide(prs, "Karbon", carbon_bullets, "esg.png", MODULE_COLORS["Karbon"])

    # Eşleştirme (Mapping)
    mapping_bullets = [
        "Amaç: SDG-GRI-TSRS eşleştirmelerini yönetme",
        "Sekmeler: SDG-GRI, SDG-TSRS, GRI-TSRS, CSV İşlemleri",
        "Butonlar: SDG seçimine göre filtrele, arama ve CSV import/export",
        "Raporlar: CSV export ile eşleştirme tabloları",
    ]
    add_bullet_slide(prs, "Eşleştirme", mapping_bullets, None, MODULE_COLORS["Eslesme"])

    # Raporlama (Ana uygulama menüsü)
    reporting_bullets = [
        "Amaç: SDG / SDG+GRI / SDG+GRI+TSRS raporlamaları",
        "Butonlar: SDG Raporu, SDG+GRI Raporu, SDG+GRI+TSRS Raporu",
        "Çıktılar: DOCX/Excel/PDF (modül raporlarına bağlı)",
    ]
    add_bullet_slide(prs, "Raporlama", reporting_bullets, "report.png", MODULE_COLORS["Raporlama"])

    # Yönetim
    mgmt_bullets = [
        "Amaç: Kullanıcı/Rol/Lisans/Modül kontrol",
        "Buton: Rol/İzin Matrisi (CSV) görüntüleme",
        "Güvenlik: Super-admin, lisans doğrulama, audit (ayrı modül)",
    ]
    add_bullet_slide(prs, "Yönetim", mgmt_bullets, "login.jpg", MODULE_COLORS["Yonetim"])

    # Kayıt
    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    return OUTPUT_PATH


def main() -> None:
    out = build_presentation()
    logging.info(f"Sunum oluşturuldu: {out}")


if __name__ == "__main__":
    main()
